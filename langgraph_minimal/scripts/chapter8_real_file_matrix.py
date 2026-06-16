from __future__ import annotations

import json
from pathlib import Path
from typing import Any

import numpy as np
import soundfile as sf
from PIL import Image, ImageDraw

from scripts.chapter8_memory_rag_dashboard_demo import build_persistent_dashboard_demo


MATRIX_MARKERS = {
    ".pdf": "PDF-MATRIX-001 PDF import test",
    ".docx": "DOCX-MATRIX-001 Word import test",
    ".pptx": "PPTX-MATRIX-001 slide import test",
    ".xlsx": "XLSX-MATRIX-001 sheet import test",
    ".png": "PNG-MATRIX-001 image OCR import test",
    ".mp3": "MP3-MATRIX-001 audio ASR import test",
    ".wav": "WAV-MATRIX-001 audio ASR import test",
}


def run_real_file_matrix(*, workspace: str | Path) -> dict[str, Any]:
    """Run the Chapter 8 real-format ingestion matrix.

    The matrix intentionally uses actual container/file formats instead of
    renaming text files.  Heavy perception models are still injected so the test
    checks the dashboard/RAG/memory workflow without loading OCR/ASR models.
    """

    root = Path(workspace)
    files_dir = root / "matrix_files"
    data_dir = root / "dashboard_data"
    files_dir.mkdir(parents=True, exist_ok=True)

    paths = _write_matrix_files(files_dir)
    dashboard = build_persistent_dashboard_demo(
        data_dir=data_dir,
        prefer_real_llm=False,
        prefer_real_multimodal=True,
        image_ocr=lambda path: MATRIX_MARKERS[".png"],
        audio_asr=lambda path: MATRIX_MARKERS[Path(path).suffix.lower()],
    )

    rows: list[dict[str, Any]] = []
    for suffix, path in paths.items():
        marker = MATRIX_MARKERS[suffix]
        if suffix in {".png", ".mp3", ".wav"}:
            import_output = dashboard.ingest_file(
                path,
                description=f"Real {suffix} matrix file.",
                importance=0.8,
            )
            modality = "image" if suffix == ".png" else "audio"
            document_id = f"perceptual:{modality}:{path.name}"
        else:
            import_output = dashboard.load_document(path)
            document_id = path.name

        search_output = dashboard.rag_tool.run(
            {
                "action": "search",
                "query": marker,
                "namespace": dashboard.rag_namespace,
                "limit": 3,
                "enable_keyword_rerank": True,
                "candidate_pool_size": 20,
            }
        )
        answer_output = dashboard.ask(f"{marker} 对应的文件说了什么？", limit=3)
        chunks = dashboard.rag_tool.document_store.list_chunks(
            document_id,
            namespace=dashboard.rag_namespace,
        )
        rows.append(
            {
                "suffix": suffix,
                "path": str(path),
                "document_id": document_id,
                "marker": marker,
                "imported": "已添加" in import_output or "Perceptual memory saved" in import_output,
                "chunk_count": len(chunks),
                "search": search_output,
                "answer": answer_output,
            }
        )

    return {
        "backend": dashboard.backend_status(),
        "files": rows,
    }


def _write_matrix_files(directory: Path) -> dict[str, Path]:
    paths = {
        ".pdf": directory / "matrix_pdf.pdf",
        ".docx": directory / "matrix_word.docx",
        ".pptx": directory / "matrix_slide.pptx",
        ".xlsx": directory / "matrix_sheet.xlsx",
        ".png": directory / "matrix_image.png",
        ".mp3": directory / "matrix_audio.mp3",
        ".wav": directory / "matrix_audio.wav",
    }
    _write_pdf(paths[".pdf"], MATRIX_MARKERS[".pdf"])
    _write_docx(paths[".docx"], MATRIX_MARKERS[".docx"])
    _write_pptx(paths[".pptx"], MATRIX_MARKERS[".pptx"])
    _write_xlsx(paths[".xlsx"], MATRIX_MARKERS[".xlsx"])
    _write_png(paths[".png"], MATRIX_MARKERS[".png"])
    _write_audio(paths[".wav"], format_name="WAV")
    _write_audio(paths[".mp3"], format_name="MP3")
    return paths


def _write_pdf(path: Path, text: str) -> None:
    stream = f"BT /F1 12 Tf 50 100 Td ({text}) Tj ET".encode("ascii")
    objects = [
        b"1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj\n",
        b"2 0 obj<</Type/Pages/Count 1/Kids[3 0 R]>>endobj\n",
        b"3 0 obj<</Type/Page/Parent 2 0 R/MediaBox[0 0 420 160]/Contents 4 0 R/Resources<</Font<</F1 5 0 R>>>>>>endobj\n",
        b"4 0 obj<</Length " + str(len(stream)).encode("ascii") + b">>stream\n" + stream + b"\nendstream endobj\n",
        b"5 0 obj<</Type/Font/Subtype/Type1/BaseFont/Helvetica>>endobj\n",
    ]
    content = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for obj in objects:
        offsets.append(len(content))
        content.extend(obj)
    xref_offset = len(content)
    content.extend(f"xref\n0 {len(offsets)}\n0000000000 65535 f \n".encode("ascii"))
    for offset in offsets[1:]:
        content.extend(f"{offset:010d} 00000 n \n".encode("ascii"))
    content.extend(
        f"trailer<</Root 1 0 R/Size {len(offsets)}>>\nstartxref\n{xref_offset}\n%%EOF".encode("ascii")
    )
    path.write_bytes(bytes(content))


def _write_docx(path: Path, text: str) -> None:
    from docx import Document

    document = Document()
    document.add_paragraph(text)
    document.save(path)


def _write_pptx(path: Path, text: str) -> None:
    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = text
    presentation.save(path)


def _write_xlsx(path: Path, text: str) -> None:
    from openpyxl import Workbook

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Matrix"
    sheet["A1"] = text
    workbook.save(path)


def _write_png(path: Path, text: str) -> None:
    image = Image.new("RGB", (900, 220), "white")
    draw = ImageDraw.Draw(image)
    draw.text((24, 90), text, fill="black")
    image.save(path)


def _write_audio(path: Path, *, format_name: str) -> None:
    samples = np.zeros(16000, dtype="float32")
    sf.write(path, samples, 16000, format=format_name)


def main() -> None:
    import argparse

    parser = argparse.ArgumentParser(description="Run Chapter 8 real file matrix.")
    parser.add_argument("--workspace", default="memory_data/real_file_matrix")
    args = parser.parse_args()
    print(json.dumps(run_real_file_matrix(workspace=args.workspace), ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
